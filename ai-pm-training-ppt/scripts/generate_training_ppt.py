# -*- coding: utf-8 -*-
"""
AI赋能项目经理实现卓越项目管理 - 培训PPT生成脚本

基于培训大纲和待裁剪PPT资源，选取并编排幻灯片，生成符合3小时培训的PPT。
用法：python generate_training_ppt.py

注意：采用「删除策略」而非「复制策略」，避免 deepcopy 导致 rId 引用断裂、PPT 损坏。
"""

import os
import shutil
import sys

# 要保留的幻灯片索引（0-based，去重后升序），与培训大纲结构一致
# 大纲：模块一(提示词+AI工具) | 模块二(项目全生命周期:启动/规划/执行监控/收尾) | 模块三(趋势)
KEEP_SLIDE_INDICES = sorted(set([
    # === 开场 ===
    0, 2, 3, 4,
    # === 模块一：提示词工程与AI工具（含自动生成项目管理计划、需求规格、设计文档研讨）===
    188, 189, 193, 196, 197, 198, 199, 200, 201, 202,
    204, 205, 206, 207, 210, 211, 214, 215, 216, 217,
    233, 234, 235, 236, 237,
    224, 225, 226, 227, 228,
    241, 242, 243, 244, 245, 246, 247,
    252, 253, 254, 255, 256, 268, 269, 270,
    # === 模块二：项目全生命周期 ===
    # 项目启动：商业价值分析(TRACE 235-237)、风险识别、干系人
    # 项目规划：需求规格(353-354)、概要设计(250-251)、WBS(252-256)、进度(452)、风险管理(437-451)
    249, 250, 251,
    353, 354,
    437, 438, 439, 440, 441, 442, 443, 444, 445, 446, 447,
    448, 449, 450, 451, 452,
    # 执行与监控：周报/会议纪要、工作流
    281, 282, 283,
    505, 506, 507, 508, 509, 510, 511, 514, 515, 516,
    # 项目收尾：交付物清单(224-227)、经验教训(RAG 186-187)
    186, 187, 204,
    # === 模块三：未来趋势 ===
    259, 330, 331, 332,
    # PMBOK框架、商业论证(453含商业分析/商业论证)
    453, 454,
]))


def create_custom_slide(prs, title, content_lines, layout_idx=16):
    """创建自定义内容页"""
    from pptx.util import Inches, Pt
    
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    
    left = Inches(0.5)
    top = Inches(0.8)
    width = Inches(12.333)
    height = Inches(6)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "微软雅黑"
    
    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "微软雅黑"
        p.space_after = Pt(6)
    
    return slide


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    skill_dir = os.path.dirname(script_dir)
    project_dir = os.path.dirname(os.path.dirname(os.path.dirname(skill_dir)))
    
    template_path = os.path.join(project_dir, "待裁剪PPT资源.pptx")
    output_path = os.path.join(project_dir, "AI赋能项目经理实现卓越项目管理培训.pptx")
    
    if not os.path.exists(template_path):
        print(f"未找到模板: {template_path}")
        sys.exit(1)
    
    from pptx import Presentation
    
    try:
        print(f"复制模板到输出: {output_path}")
        shutil.copy2(template_path, output_path)
    except PermissionError:
        alt_path = os.path.join(project_dir, "AI赋能项目经理实现卓越项目管理培训_新稿.pptx")
        print(f"原文件可能被占用，改用: {alt_path}")
        output_path = alt_path
        shutil.copy2(template_path, output_path)
    
    prs = Presentation(output_path)
    total = len(prs.slides)
    keep_set = set(KEEP_SLIDE_INDICES)
    
    # 从后往前删除不需要的幻灯片（避免索引错位）
    removed = 0
    sldIdLst = prs.slides._sldIdLst
    slides_list = list(sldIdLst)
    for i in range(total - 1, -1, -1):
        if i not in keep_set:
            try:
                elem = slides_list[i]
                rId = elem.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                if rId:
                    prs.part.drop_rel(rId)
                elem.getparent().remove(elem)
                removed += 1
                if removed % 50 == 0:
                    print(f"  已删除 {removed} 页...")
            except Exception as ex:
                print(f"  删除第{i+1}页失败: {ex}")
    
    # 修改封面标题
    if len(prs.slides) > 0:
        slide0 = prs.slides[0]
        for shp in slide0.shapes:
            if hasattr(shp, "text") and shp.text and ("大模型" in shp.text or "01" in shp.text):
                shp.text = "AI赋能项目经理实现卓越项目管理\n——XX电信三小时专题培训"
                break
    
    prs.save(output_path)
    print(f"\n生成完成: {output_path}")
    print(f"共 {len(prs.slides)} 页，约 {len(prs.slides) * 1.5:.0f} 分钟（按1.5分钟/页）")


if __name__ == "__main__":
    main()
