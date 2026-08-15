# -*- coding: utf-8 -*-
"""
从待裁剪PPT资源构建包含三个模块的完整培训PPT

按照新的目录结构：
- 模块一：提示词工程与特定AI工具在项目管理中的应用（45分钟）
- 模块二：AI在项目全生命周期的应用场景与实操演练（105分钟）
- 模块三：AI驱动的项目管理未来趋势与组织应对（15分钟）

用法：python build_ppt_with_three_modules.py
"""
import os
import shutil
import sys

# 项目根目录
PROJECT_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))))

# 根据reference.md的映射，定义三个模块的页面索引（0-based）
# 模块一：提示词工程与特定AI工具在项目管理中的应用（45分钟）
MODULE_1_INDICES = sorted(set([
    # 典型提示词技巧
    188, 189, 193, 196, 197, 198, 199, 200, 201, 202, 204, 205, 206, 207, 210, 211, 214, 215, 216, 217, 224, 225, 226, 227, 228,
    # 自动生成项目管理计划
    241, 242, 243, 244, 245, 246, 247,
    # AI工具 Mermaid/Coze/Dify
    268, 269, 270,
    # AI智能体与工作流
    505, 506, 507, 508, 509, 510, 511, 514, 515, 516,
    # 需求规格/架构设计自动化研讨
    249, 250, 251, 252, 253, 254, 255, 256, 353, 354,
]))

# 模块二：AI在项目全生命周期的应用场景与实操演练（105分钟）
MODULE_2_INDICES = sorted(set([
    # 项目启动：商业价值分析
    233, 234, 235, 236, 237, 452, 453, 454,
    # 风险识别、干系人
    437, 438, 439, 440, 441, 442, 443, 444, 445, 446, 447,
    # 政企项目实操
    238, 239, 240,
    # 项目规划：需求分析文档
    249, 250, 251, 353, 354,
    # 设计文档+WBS
    448, 449,
    # 进度估算与关键路径
    508, 509, 510, 511,
    # 风险管理计划
    # 执行与监控：周报/会议纪要
    281, 282, 283,
    # 进度/成本偏差预警（已在工作流中）
    # 项目收尾：交付物清单
    224, 225, 226, 227,
    # 经验教训知识库
    186, 187, 204,
    # 模块二过渡页
    248,
]))

# 模块三：AI驱动的项目管理未来趋势与组织应对（15分钟）
MODULE_3_INDICES = sorted(set([
    259, 330, 331, 332,
]))

# 前10页保留（开场页）
KEEP_FIRST_10_INDICES = [0, 2, 3, 4, 5, 6, 7, 8, 9, 10]

# 所有需要保留的页面
ALL_KEEP_INDICES = sorted(set(KEEP_FIRST_10_INDICES + MODULE_1_INDICES + MODULE_2_INDICES + MODULE_3_INDICES))

# 新目录（第11页）
NEW_TOC_LINES = [
    "模块一：提示词工程与特定AI工具在项目管理中的应用（45分钟）",
    "",
    "模块二：AI在项目全生命周期的应用场景与实操演练（105分钟）",
    "",
    "模块三：AI驱动的项目管理未来趋势与组织应对（15分钟）",
]

# 更详细的目录（可选）
DETAILED_TOC_LINES = [
    "模块一：提示词工程与特定AI工具在项目管理中的应用（45分钟）",
    "  典型的提示词技巧有哪些？",
    "  基于提示词框架来思考如何自动生成项目管理计划",
    "  典型的AI工具介绍：Mermaid、PlantUML、Coze、Dify等",
    "  AI智能体与工作流如何助力项目管理的流程优化？",
    "  研讨项目需求规格说明书、架构设计文档的关键组件的自动化生成的可能性",
    "",
    "模块二：AI在项目全生命周期的应用场景与实操演练（105分钟）",
    "  （1）项目启动：市场调研、风险识别、干系人识别、政企实操",
    "  （2）项目规划：需求文档、设计文档+WBS、进度估算、风险管理",
    "  （3）执行与监控：周报/会议纪要、进度/成本偏差预警",
    "  （4）项目收尾：交付物清单、经验教训知识库",
    "",
    "模块三：AI驱动的项目管理未来趋势与组织应对（15分钟）",
    "  趋势1：多个AI智能体的协同",
    "  趋势2：本地知识库+通用大模型",
    "  趋势3：PM角色进化",
    "  组织建议与互动问答",
]


def _safe_remove_slide(prs, index):
    """安全删除幻灯片"""
    try:
        sldIdLst = prs.slides._sldIdLst
        slides_list = list(sldIdLst)
        if index >= len(slides_list):
            return False
        elem = slides_list[index]
        rId = elem.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId:
            prs.part.drop_rel(rId)
        elem.getparent().remove(elem)
        return True
    except Exception as ex:
        print(f"  删除第{index+1}页失败: {ex}")
        return False


def _update_toc_slide(prs):
    """更新第11页为新的课程大纲目录"""
    from pptx.util import Pt
    if len(prs.slides) < 11:
        return
    slide = prs.slides[10]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full_text = "".join(p.text for p in tf.paragraphs).strip()
        if any(k in full_text for k in ("目录", "CATALOGUE", "大模型基础", "与大模型沟通", "大模型智能办公", "零代码", "AI 智能编程", "课程结构", "模块一", "模块二", "模块三")):
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = "目录"
            p0.font.size = Pt(32)
            p0.font.bold = True
            p0.font.name = "微软雅黑"
            for line in DETAILED_TOC_LINES:
                p = tf.add_paragraph()
                p.text = line
                if line.startswith("模块"):
                    p.font.size = Pt(22)
                    p.font.bold = True
                elif line.startswith("  "):
                    p.font.size = Pt(16)
                    p.font.bold = False
                elif line == "":
                    p.font.size = Pt(8)
                else:
                    p.font.size = Pt(18)
                p.font.name = "微软雅黑"
                p.space_before = Pt(12)
            return
    # 未找到：更新主内容区
    best = None
    for shape in slide.shapes:
        if shape.has_text_frame and len(shape.text_frame.paragraphs) > 0:
            if best is None or len(shape.text) > len(best.text):
                best = shape
    if best:
        tf = best.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = "目录"
        p0.font.size = Pt(32)
        p0.font.bold = True
        p0.font.name = "微软雅黑"
        for line in DETAILED_TOC_LINES:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("模块"):
                p.font.size = Pt(22)
                p.font.bold = True
            elif line.startswith("  "):
                p.font.size = Pt(16)
                p.font.bold = False
            elif line == "":
                p.font.size = Pt(8)
            else:
                p.font.size = Pt(18)
            p.font.name = "微软雅黑"
            p.space_before = Pt(12)


def build_from_source(template_path, output_path):
    """从待裁剪PPT资源构建课件：保留开场10页 + 三个模块对应页"""
    from pptx import Presentation

    prs = Presentation(template_path)
    total = len(prs.slides)
    
    # 要保留的页码
    keep_set = {i for i in ALL_KEEP_INDICES if i < total}
    
    print(f"待裁剪PPT总页数: {total}")
    print(f"模块一页面数: {len([i for i in MODULE_1_INDICES if i < total])}")
    print(f"模块二页面数: {len([i for i in MODULE_2_INDICES if i < total])}")
    print(f"模块三页面数: {len([i for i in MODULE_3_INDICES if i < total])}")
    print(f"将保留总页数: {len(keep_set)}")
    
    removed = 0
    for i in range(total - 1, -1, -1):
        if i not in keep_set:
            if _safe_remove_slide(prs, i):
                removed += 1
                if removed % 50 == 0:
                    print(f"  已删除 {removed} 页...")

    # 修改封面
    if len(prs.slides) > 0:
        slide0 = prs.slides[0]
        for shp in slide0.shapes:
            if hasattr(shp, "text") and shp.text:
                if "大模型" in shp.text or "01" in shp.text:
                    shp.text = "AI赋能项目经理实现卓越项目管理\n——XX电信三小时专题培训"
                    break

    # 第11页改为新课程大纲目录
    _update_toc_slide(prs)

    prs.save(output_path)
    return len(prs.slides)


def main():
    project_dir = PROJECT_DIR

    # 使用待裁剪PPT资源作为源文件
    source_path = os.path.join(project_dir, "待裁剪PPT资源.pptx")
    output_path = os.path.join(project_dir, "AI赋能项目经理实现卓越项目管理培训.pptx")

    if not os.path.exists(source_path):
        print("未找到 待裁剪PPT资源.pptx，请将文件放入项目目录：")
        print(f"  {project_dir}")
        sys.exit(1)

    print("从待裁剪PPT资源构建（包含三个模块）...")
    try:
        # 备份现有文件（如果存在）
        if os.path.exists(output_path):
            backup_path = output_path.replace(".pptx", "_备份.pptx")
            shutil.copy2(output_path, backup_path)
            print(f"已备份现有文件至: {backup_path}")
        
        shutil.copy2(source_path, output_path)
    except PermissionError:
        output_path = os.path.join(project_dir, "AI赋能项目经理实现卓越项目管理培训_新稿.pptx")
        shutil.copy2(source_path, output_path)
    
    count = build_from_source(output_path, output_path)

    print(f"\n生成完成: {output_path}")
    print(f"共 {count} 页，约 {count * 1.5:.0f} 分钟（按1.5分钟/页）")
    print(f"\n目录结构：")
    print(f"  模块一：提示词工程与特定AI工具在项目管理中的应用（45分钟）")
    print(f"  模块二：AI在项目全生命周期的应用场景与实操演练（105分钟）")
    print(f"  模块三：AI驱动的项目管理未来趋势与组织应对（15分钟）")


if __name__ == "__main__":
    main()
