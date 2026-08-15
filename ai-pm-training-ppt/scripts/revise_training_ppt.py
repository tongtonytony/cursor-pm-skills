# -*- coding: utf-8 -*-
"""
AI赋能项目经理实现卓越项目管理 - 培训PPT修订脚本

基于课程大纲修订PPT：保留第1-10页，从第11页起按大纲结构重新编排内容。
用法：python revise_training_ppt.py

策略：仅采用「删除策略」，从待裁剪PPT中保留所需页并删除其余页。
绝不跨文件复制幻灯片，避免 rId 引用断裂导致的「PowerPoint 无法读取部分内容」报错。
"""
import os
import shutil
import sys

# 项目根目录（scripts -> ai-pm-training-ppt -> skills -> .cursor -> 项目管理）
PROJECT_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))))

# 与《AI赋能项目经理实现卓越项目管理培训大纲.docx》完全一致的目录结构
# 格式：(模块/章节名, [待裁剪PPT页码列表 0-based])
# 涵盖：商业论证、项目章程、需求规格说明书、概要设计、WBS、进度网络图、项目管理计划
OUTLINE_SLIDE_MAP = [
    # ========== 第11页起 ==========
    # 模块一：提示词工程与特定AI工具（45分钟）
    ("模块一：提示词工程与特定AI工具", [188, 189]),
    ("典型的提示词技巧", [193, 196, 197, 198, 199, 200, 201, 202, 204, 205, 206, 207, 210, 211, 214, 215]),
    ("提示词框架（TRACE/ICIO/CRISPE等）", [216, 217, 228, 229, 230, 231, 232, 233, 234]),
    ("基于提示词框架自动生成项目管理计划", [241, 242, 243, 244, 245, 246, 247]),
    ("典型AI工具：Mermaid、PlantUML、Coze、Dify", [268, 269, 270]),
    ("AI智能体与工作流助力项目管理流程优化", [505, 506, 507, 508, 509, 510, 511, 514, 515, 516]),
    ("研讨：需求规格说明书、架构设计文档自动化生成", [245, 246, 247, 249, 250, 251, 252, 253, 254, 255, 256, 353, 354]),
    # ========== 模块二：AI在项目全生命周期的应用（105分钟）==========
    ("模块二：AI在项目全生命周期的应用", [248]),
    # 2.1 项目启动：商业论证、项目章程
    ("项目启动：AI辅助市场调研与商业价值分析", [233, 234, 235, 236, 237]),
    ("商业论证与项目章程自动生成", [452, 453, 454]),
    ("AI辅助识别项目初始风险、干系人识别", [437, 438, 439, 440, 441, 442, 443, 444, 445, 446, 447]),
    ("实操：用典型政企项目举例", [235, 236, 237, 238, 239, 240]),
    # 2.2 项目规划：需求、设计、WBS、进度、风险管理
    ("使用AI编写需求分析文档、需求规格说明书", [249, 250, 251, 353, 354]),
    ("使用AI编写设计文档、项目概要设计说明书", [250, 251]),
    ("WBS智能分解与范围管理", [252, 253, 254, 255, 256, 448, 449]),
    ("进度估算与关键路径、项目进度网络图", [452, 453, 508, 509, 510, 511]),
    ("项目管理计划与全生命周期文档自动化", [452, 453, 454]),
    ("风险管理计划：合同/需求/设计文档自动提取风险", [437, 438, 439, 440, 441, 442, 443, 444, 445, 446, 447]),
    # 2.3 执行与监控
    ("自动撰写项目周报/会议纪要（语音转写+摘要）", [281, 282, 283]),
    ("进度偏差预警、成本偏差预警", [505, 506, 507, 508, 509, 510, 511]),
    # 2.4 项目收尾
    ("项目交付物清单核对跟踪表自动生成", [224, 225, 226, 227]),
    ("经验教训知识库自动生成与归档", [186, 187, 204]),
    # ========== 模块三：AI驱动的项目管理未来趋势（15分钟）==========
    ("模块三：AI驱动的项目管理未来趋势与组织应对", [259]),
    ("未来趋势、组织建议、互动问答（佛山电信）", [259, 330, 331, 332]),
]

# 按大纲顺序去重后的幻灯片索引（用于「从待裁剪PPT删除策略」）
def _flatten_outline_indices():
    seen = set()
    result = []
    for _, indices in OUTLINE_SLIDE_MAP:
        for i in indices:
            if i not in seen:
                seen.add(i)
                result.append(i)
    return result

OUTLINE_SLIDE_INDICES = _flatten_outline_indices()

# 前10页保留的页码（0-based）。当基稿为培训PPT时保留；当从待裁剪生成时，用此作为开场
# 待裁剪PPT中：0=大模型基础, 2=多模态, 3-4=章节重点, 可扩展为10页
KEEP_FIRST_10_INDICES = [0, 2, 3, 4, 5, 6, 7, 8, 9, 10]  # 开场约10页


def _safe_remove_slide(prs, index):
    """安全删除幻灯片，避免 rId 残留导致损坏"""
    try:
        sldIdLst = prs.slides._sldIdLst
        slides_list = list(sldIdLst)
        if index >= len(slides_list):
            return False
        elem = slides_list[index]
        rId = elem.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId:
            prs.part.drop_rel(rId)
        elem.getparent().remove(elem)
        return True
    except Exception as ex:
        print(f"  删除第{index+1}页失败: {ex}")
        return False


# 课程大纲目录（与培训大纲完全一致，用于第11页）
TOC_LINES = [
    "1. 模块一：提示词工程与特定AI工具在项目管理中的应用（45分钟）",
    "2. 模块二：AI在项目全生命周期的应用场景与实操演练（105分钟）",
    "3. 模块三：AI驱动的项目管理未来趋势与组织应对（15分钟）",
]


def _update_toc_slide(prs):
    """将第11页（0-based索引10）内容更新为课程大纲目录"""
    from pptx.util import Pt
    if len(prs.slides) < 11:
        return
    slide = prs.slides[10]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full_text = "".join(p.text for p in tf.paragraphs).strip()
        # 匹配旧目录（大模型基础、与大模型沟通等）或已有目录
        if any(k in full_text for k in ("目录", "CATALOGUE", "大模型基础", "与大模型沟通", "大模型智能办公", "零代码", "AI 智能编程")):
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = "目录"
            p0.font.size = Pt(28)
            p0.font.bold = True
            p0.font.name = "微软雅黑"
            for line in TOC_LINES:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(18)
                p.font.name = "微软雅黑"
                p.space_before = Pt(12)
            return
    # 未找到旧目录：更新主内容区（通常为最大文本块）
    best = None
    for shape in slide.shapes:
        if shape.has_text_frame and len(shape.text_frame.paragraphs) > 0:
            if best is None or len(shape.text) > len(best.text):
                best = shape
    if best:
        tf = best.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = "目录"
        p0.font.size = Pt(28)
        p0.font.bold = True
        p0.font.name = "微软雅黑"
        for line in TOC_LINES:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.name = "微软雅黑"
            p.space_before = Pt(12)


def build_from_source(template_path, output_path):
    """
    从待裁剪PPT资源构建课件：保留开场10页 + 大纲对应页。
    采用删除策略，从后往前删除。
    """
    from pptx import Presentation

    prs = Presentation(template_path)
    total = len(prs.slides)
    # 要保留的页码 = 前10 + 大纲页
    keep_set = set(KEEP_FIRST_10_INDICES) | set(OUTLINE_SLIDE_INDICES)
    # 确保不超出范围
    keep_set = {i for i in keep_set if i < total}
    removed = 0
    for i in range(total - 1, -1, -1):
        if i not in keep_set:
            if _safe_remove_slide(prs, i):
                removed += 1
                if removed % 50 == 0:
                    print(f"  已删除 {removed} 页...")

    # 修改封面
    if len(prs.slides) > 0:
        slide0 = prs.slides[0]
        for shp in slide0.shapes:
            if hasattr(shp, "text") and shp.text:
                if "大模型" in shp.text or "01" in shp.text:
                    shp.text = "AI赋能项目经理实现卓越项目管理\n——XX电信三小时专题培训"
                    break

    # 第11页改为课程大纲目录（与培训大纲完全一致）
    _update_toc_slide(prs)

    prs.save(output_path)
    return len(prs.slides)


def main():
    project_dir = PROJECT_DIR

    source_path = os.path.join(project_dir, "待裁剪PPT资源.pptx")
    output_path = os.path.join(project_dir, "AI赋能项目经理实现卓越项目管理培训.pptx")

    if not os.path.exists(source_path):
        print("未找到 待裁剪PPT资源.pptx，请将文件放入项目目录：")
        print(f"  {project_dir}")
        sys.exit(1)

    # 仅使用删除策略：复制待裁剪PPT后删除不需页，避免跨文件复制导致的 rId 损坏
    print("从待裁剪PPT资源构建（纯删除策略，避免PPT打开报错）...")
    try:
        shutil.copy2(source_path, output_path)
    except PermissionError:
        output_path = os.path.join(project_dir, "AI赋能项目经理实现卓越项目管理培训_新稿.pptx")
        shutil.copy2(source_path, output_path)
    count = build_from_source(output_path, output_path)

    print(f"\n生成完成: {output_path}")
    print(f"共 {count} 页，约 {count * 1.5:.0f} 分钟（按1.5分钟/页）")


if __name__ == "__main__":
    main()
