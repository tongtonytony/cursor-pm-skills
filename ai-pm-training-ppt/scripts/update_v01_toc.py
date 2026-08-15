# -*- coding: utf-8 -*-
"""
更新 AI赋能项目经理实现卓越项目管理培训V0.1.pptx 的第11页目录

用法：python update_v01_toc.py
"""
import os
import shutil
import sys

# 项目根目录
PROJECT_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))))

# 新目录内容（第11页）
TOC_LINES = [
    "模块一：提示词工程与特定AI工具在项目管理中的应用（45分钟）",
    "典型的提示词技巧有哪些？",
    "基于提示词框架来思考如何自动生成项目管理计划",
    "典型的AI工具介绍：Mermaid、PlantUML、Coze、Dify等",
    "AI智能体与工作流如何助力项目管理的流程优化？",
    "研讨项目需求规格说明书、架构设计文档的关键组件的自动化生成的可能性",
    "",
    "模块二：AI在项目全生命周期的应用场景与实操演练（105分钟）",
    "（1）项目启动：市场调研、风险识别、干系人识别、政企实操",
    "（2）项目规划：需求文档、设计文档+WBS、进度估算、风险管理",
    "（3）执行与监控：周报/会议纪要、进度/成本偏差预警",
    "（4）项目收尾：交付物清单、经验教训知识库",
    "",
    "模块三：AI驱动的项目管理未来趋势与组织应对（15分钟）",
    "趋势1：多个AI智能体的协同",
    "趋势2：本地知识库+通用大模型",
    "趋势3：PM角色进化",
]


def update_toc_slide(prs):
    """更新第11页（索引10）为新的课程大纲目录"""
    from pptx.util import Pt
    
    if len(prs.slides) < 11:
        print(f"警告：PPT只有 {len(prs.slides)} 页，无法更新第11页")
        return False
    
    slide = prs.slides[10]  # 第11页，索引为10
    
    # 查找并更新文本框架
    updated = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        tf = shape.text_frame
        full_text = "".join(p.text for p in tf.paragraphs).strip()
        
        # 如果找到目录相关的内容，或者找到最大的文本框架，则更新
        if any(k in full_text for k in ("目录", "CATALOGUE", "模块一", "模块二", "模块三", "课程结构")) or len(full_text) > 20:
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = "目录"
            p0.font.size = Pt(32)
            p0.font.bold = True
            p0.font.name = "微软雅黑"
            
            for line in TOC_LINES:
                p = tf.add_paragraph()
                p.text = line
                
                # 设置字体大小和样式
                if line.startswith("模块"):
                    p.font.size = Pt(22)
                    p.font.bold = True
                    p.space_before = Pt(12)
                elif line.startswith("（") or line.startswith("趋势"):
                    p.font.size = Pt(18)
                    p.font.bold = False
                    p.space_before = Pt(6)
                elif line == "":
                    p.font.size = Pt(8)
                    p.space_before = Pt(4)
                else:
                    p.font.size = Pt(18)
                    p.font.bold = False
                    p.space_before = Pt(6)
                
                p.font.name = "微软雅黑"
            
            updated = True
            break
    
    # 如果没找到合适的文本框架，尝试更新最大的文本框架
    if not updated:
        best = None
        max_len = 0
        for shape in slide.shapes:
            if shape.has_text_frame and len(shape.text_frame.paragraphs) > 0:
                text_len = len(shape.text)
                if text_len > max_len:
                    max_len = text_len
                    best = shape
        
        if best:
            tf = best.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = "目录"
            p0.font.size = Pt(32)
            p0.font.bold = True
            p0.font.name = "微软雅黑"
            
            for line in TOC_LINES:
                p = tf.add_paragraph()
                p.text = line
                
                if line.startswith("模块"):
                    p.font.size = Pt(22)
                    p.font.bold = True
                    p.space_before = Pt(12)
                elif line.startswith("（") or line.startswith("趋势"):
                    p.font.size = Pt(18)
                    p.font.bold = False
                    p.space_before = Pt(6)
                elif line == "":
                    p.font.size = Pt(8)
                    p.space_before = Pt(4)
                else:
                    p.font.size = Pt(18)
                    p.font.bold = False
                    p.space_before = Pt(6)
                
                p.font.name = "微软雅黑"
            updated = True
    
    return updated


def main():
    project_dir = PROJECT_DIR
    
    # 目标文件
    ppt_path = os.path.join(project_dir, "AI赋能项目经理实现卓越项目管理培训V0.1.pptx")
    
    if not os.path.exists(ppt_path):
        print(f"未找到文件: {ppt_path}")
        print(f"请确认文件路径: {project_dir}")
        sys.exit(1)
    
    try:
        from pptx import Presentation
        
        print(f"读取PPT文件: {ppt_path}")
        prs = Presentation(ppt_path)
        
        print(f"PPT总页数: {len(prs.slides)}")
        
        # 备份原文件
        backup_path = ppt_path.replace(".pptx", "_备份.pptx")
        shutil.copy2(ppt_path, backup_path)
        print(f"已备份原文件至: {backup_path}")
        
        # 更新目录页
        if update_toc_slide(prs):
            prs.save(ppt_path)
            print(f"\n成功更新第11页目录")
            print(f"保存文件: {ppt_path}")
        else:
            print("\n未能找到合适的目录页进行更新")
            sys.exit(1)
        
    except Exception as ex:
        print(f"错误: {ex}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
