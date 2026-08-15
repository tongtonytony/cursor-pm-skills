# -*- coding: utf-8 -*-
"""
AI赋能项目经理培训PPT - 源文件内容提取脚本

用法：
  python extract_source_content.py <大纲docx路径> <待裁剪PPT路径>

输出：
  - outline_extracted.txt：培训大纲章节结构及要点
  - pptx_extracted.txt：待裁剪PPT每页文本及布局说明
"""

import os
import sys


def extract_docx(path: str) -> str:
    """提取 docx 全文文本（保留段落结构）"""
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        return "[需要安装 python-docx: pip install python-docx]"

    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # 根据样式粗略判断层级
            style = para.style.name if para.style else ""
            if "Heading" in style or "标题" in style:
                lines.append("\n" + text + "\n")
            else:
                lines.append(text)
    return "\n".join(lines).strip()


def extract_pptx(path: str) -> str:
    """提取 pptx 每页文本及布局说明"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "[需要安装 python-pptx: pip install python-pptx]"

    prs = Presentation(path)
    pages = []

    for idx, slide in enumerate(prs.slides, start=1):
        page_lines = [f"========== 第 {idx} 页 =========="]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                page_lines.append(shape.text.strip())
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        page_lines.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        page_lines.append(row_text)

        if len(page_lines) == 1:
            page_lines.append("(无文本内容)")

        pages.append("\n".join(page_lines))

    return "\n\n".join(pages)


def main():
    if len(sys.argv) < 3:
        print("用法: python extract_source_content.py <大纲docx路径> <待裁剪PPT路径>")
        print("示例: python extract_source_content.py AI赋能项目经理实现卓越项目管理培训大纲.docx 待裁剪PPT资源.pptx")
        sys.exit(1)

    outline_path = sys.argv[1]
    pptx_path = sys.argv[2]
    out_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    if not os.path.exists(outline_path):
        print(f"未找到大纲文件: {outline_path}")
        sys.exit(1)
    if not os.path.exists(pptx_path):
        print(f"未找到PPT文件: {pptx_path}")
        sys.exit(1)

    print("提取培训大纲...")
    outline_text = extract_docx(outline_path)
    outline_out = os.path.join(out_dir, "outline_extracted.txt")
    with open(outline_out, "w", encoding="utf-8") as f:
        f.write(outline_text)
    print(f"  -> {outline_out}")

    print("提取待裁剪PPT...")
    pptx_text = extract_pptx(pptx_path)
    pptx_out = os.path.join(out_dir, "pptx_extracted.txt")
    with open(pptx_out, "w", encoding="utf-8") as f:
        f.write(pptx_text)
    print(f"  -> {pptx_out}")

    print("提取完成。请将两个输出文件内容提供给大模型进行PPT设计。")

if __name__ == "__main__":
    main()
