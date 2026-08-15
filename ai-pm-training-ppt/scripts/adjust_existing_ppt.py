# -*- coding: utf-8 -*-
"""
基于现有PPT，按新目录结构调整

保持现有PPT的格式和内容，仅调整目录页和页面顺序。
用法：python adjust_existing_ppt.py
"""
import os
import shutil
import sys

# 项目根目录
PROJECT_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))))

# 新目录（第11页）
NEW_TOC_LINES = [
    "模块一：提示词工程与特定AI工具在项目管理中的应用（45分钟）",
    "  典型的提示词技巧有哪些？",
    "  基于提示词框架来思考如何自动生成项目管理计划",
    "  典型的AI工具介绍：Mermaid、PlantUML、Coze、Dify等",
    "  AI智能体与工作流如何助力项目管理的流程优化？",
    "  研讨项目需求规格说明书、架构设计文档的关键组件的自动化生成的可能性",
    "",
    "模块二：AI在项目全生命周期的应用场景与实操演练（105分钟）",
    "  （1）项目启动：市场调研、风险识别、干系人识别、政企实操",
    "  （2）项目规划：需求文档、设计文档+WBS、进度估算、风险管理",
    "  （3）执行与监控：周报/会议纪要、进度/成本偏差预警",
    "  （4）项目收尾：交付物清单、经验教训知识库",
    "",
    "模块三：AI驱动的项目管理未来趋势与组织应对（15分钟）",
    "  趋势1：多个AI智能体的协同",
    "  趋势2：本地知识库+通用大模型",
    "  趋势3：PM角色进化",
    "  组织建议与互动问答",
]


def _update_toc_slide(prs):
    """更新第11页为新的课程大纲目录"""
    from pptx.util import Pt
    if len(prs.slides) < 11:
        return
    slide = prs.slides[10]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full_text = "".join(p.text for p in tf.paragraphs).strip()
        if any(k in full_text for k in ("目录", "CATALOGUE", "课程结构", "模块一", "模块二", "模块三")):
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = "目录"
            p0.font.size = Pt(28)
            p0.font.bold = True
            p0.font.name = "微软雅黑"
            for line in NEW_TOC_LINES:
                p = tf.add_paragraph()
                p.text = line
                if line.startswith("模块"):
                    p.font.size = Pt(20)
                    p.font.bold = True
                elif line.startswith("  "):
                    p.font.size = Pt(16)
                    p.font.bold = False
                elif line == "":
                    p.font.size = Pt(8)
                else:
                    p.font.size = Pt(18)
                p.font.name = "微软雅黑"
                p.space_before = Pt(6)
            return
    # 未找到：更新主内容区
    best = None
    for shape in slide.shapes:
        if shape.has_text_frame and len(shape.text_frame.paragraphs) > 0:
            if best is None or len(shape.text) > len(best.text):
                best = shape
    if best:
        tf = best.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = "目录"
        p0.font.size = Pt(28)
        p0.font.bold = True
        p0.font.name = "微软雅黑"
        for line in NEW_TOC_LINES:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("模块"):
                p.font.size = Pt(20)
                p.font.bold = True
            elif line.startswith("  "):
                p.font.size = Pt(16)
                p.font.bold = False
            elif line == "":
                p.font.size = Pt(8)
            else:
                p.font.size = Pt(18)
            p.font.name = "微软雅黑"
            p.space_before = Pt(6)


def main():
    project_dir = PROJECT_DIR

    # 使用现有的PPT文件
    source_path = os.path.join(project_dir, "AI赋能项目经理实现卓越项目管理培训.pptx")
    output_path = os.path.join(project_dir, "AI赋能项目经理实现卓越项目管理培训.pptx")

    if not os.path.exists(source_path):
        print(f"未找到现有PPT文件: {source_path}")
        sys.exit(1)

    from pptx import Presentation

    print("读取现有PPT并更新目录页...")
    try:
        # 备份原文件
        backup_path = source_path.replace(".pptx", "_备份.pptx")
        shutil.copy2(source_path, backup_path)
        print(f"已备份原文件至: {backup_path}")
        
        prs = Presentation(source_path)
        
        # 更新目录页
        _update_toc_slide(prs)
        
        prs.save(output_path)
        print(f"\n调整完成: {output_path}")
        print(f"共 {len(prs.slides)} 页")
        
    except Exception as ex:
        print(f"错误: {ex}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
