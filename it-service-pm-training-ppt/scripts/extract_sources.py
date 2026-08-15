# -*- coding: utf-8 -*-
"""Extract outline docx, baseline pptx, and optional source pool summaries."""

import glob
import json
import os
import sys


def extract_docx(path: str) -> str:
    from docx import Document

    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if "Heading" in style or "标题" in style:
            lines.append("\n" + text + "\n")
        else:
            lines.append(text)
    return "\n".join(lines).strip()


def extract_pptx(path: str) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    pages = []
    for idx, slide in enumerate(prs.slides, start=1):
        page_lines = [f"========== 第 {idx} 页 =========="]
        pic_count = sum(
            1
            for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE and s.width > 800000
        )
        page_lines.append(f"[shapes={len(slide.shapes)} pics={pic_count}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                page_lines.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        page_lines.append(row_text)
        if len(page_lines) == 2:
            page_lines.append("(无文本内容)")
        pages.append("\n".join(page_lines))
    return "\n\n".join(pages)


def main():
    if len(sys.argv) < 4:
        print("用法: python extract_sources.py <大纲docx> <基线pptx> <工作目录>")
        sys.exit(1)

    outline_path, baseline_path, work_dir = sys.argv[1:4]
    os.makedirs(work_dir, exist_ok=True)

    outline_text = extract_docx(outline_path)
    baseline_text = extract_pptx(baseline_path)

    outline_out = os.path.join(work_dir, "outline_extracted.txt")
    baseline_out = os.path.join(work_dir, "baseline_extracted.txt")
    with open(outline_out, "w", encoding="utf-8") as f:
        f.write(outline_text)
    with open(baseline_out, "w", encoding="utf-8") as f:
        f.write(baseline_text)

    # optional outline json in work dir
    outline_json = os.path.join(work_dir, "_outline_chapters.json")
    if os.path.exists(outline_json):
        with open(outline_json, encoding="utf-8") as f:
            chapters = json.load(f)
        chap_out = os.path.join(work_dir, "outline_chapters_summary.txt")
        with open(chap_out, "w", encoding="utf-8") as f:
            for i, ch in enumerate(chapters, 1):
                f.write(f"\n=== 章节 {i} ===\n{ch.get('module','')}\n{ch.get('topics','')}\n")

    print(f"Wrote {outline_out}")
    print(f"Wrote {baseline_out}")
    if os.path.exists(outline_json):
        print(f"Wrote {chap_out}")


if __name__ == "__main__":
    main()
