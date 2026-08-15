# -*- coding: utf-8 -*-
"""
Template-based slide insert + PREP notes enrichment.
Phases: insert | notes | fix-dividers | all
"""

import argparse
import copy
import glob
import io
import json
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Pt

REF_LEFT, REF_TOP = 1409700, 152400
TOLERANCE = 500000
TITLE_FONT = "Microsoft YaHei"
ACCENT = RGBColor(0x00, 0x6B, 0xA6)
DARK = RGBColor(0x1A, 0x1A, 0x1A)


def norm(s):
    s = re.sub(r"---此页放到.*", "", s or "")
    s = re.sub(r"（插）", "", s)
    s = re.sub(r"\s+", "", s)
    s = re.sub(r"[：:—\-–—|｜/\\（）()【】\[\]《》…·?？]", "", s)
    return s.lower()


def get_title_shape(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if abs(shape.left - REF_LEFT) < TOLERANCE and abs(shape.top - REF_TOP) < TOLERANCE:
            if shape.text_frame.text.strip():
                return shape
    for shape in sorted(
        (s for s in slide.shapes if s.has_text_frame and s.top < 900000),
        key=lambda s: s.top,
    ):
        t = shape.text_frame.text.strip()
        if t and len(t.split("\n")[0]) < 100:
            return shape
    return None


def get_title(slide):
    shape = get_title_shape(slide)
    return shape.text_frame.text.strip().split("\n")[0] if shape else ""


def clean_title(text):
    return re.sub(r"---此页放到「.+?」后面", "", text or "").strip()


def set_run_font(run, size_pt, bold=False, color=DARK):
    run.font.name = TITLE_FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def set_title_style(shape, text):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        set_run_font(run, 30, bold=False)


def set_body_text(shape, text, size_pt=18):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(4)
        for run in p.runs:
            set_run_font(run, size_pt, bold=False)


def is_logo_pic(shape):
    return shape.left > 9000000 and shape.top < 500000


def is_decorative(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and is_logo_pic(shape):
        return True
    if shape.top < 800000 and shape.left < 500000 and shape.has_text_frame:
        if not shape.text_frame.text.strip():
            return True
    return False


def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    el = lst[from_idx]
    lst.remove(el)
    lst.insert(to_idx, el)


def duplicate_slide(prs, index):
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split("}")[-1]
        if tag in ("sp", "pic", "graphicFrame", "cxnSp", "grpSp"):
            sp_tree.remove(child)
    for shape in source.shapes:
        sp_tree.insert_element_before(copy.deepcopy(shape.element), "p:extLst")
    return new_slide


def find_template_index(prs, marker="XXXX（插）"):
    for i, slide in enumerate(prs.slides):
        t = get_title(slide)
        if "XXXX" in t or t == marker:
            return i
    for i, slide in enumerate(prs.slides):
        if get_title(slide).endswith("（插）") and len(slide.shapes) <= 6:
            return i
    return min(26, len(prs.slides) - 1)


def add_blue_divider(slide):
    for shape in slide.shapes:
        if (
            shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and shape.width > 5000000
            and shape.height < 80000
            and 500000 < shape.top < 800000
        ):
            return
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, REF_LEFT, 600000, 7696200, 28575)
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def clear_content_shapes(slide, title_shape):
    for shape in list(slide.shapes):
        if shape is title_shape or is_decorative(shape):
            continue
        if shape.has_text_frame and shape.top < 900000 and shape is not title_shape:
            continue
        tag = shape.shape_type
        if tag in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX):
            if is_decorative(shape):
                continue
            if shape.has_text_frame and abs(shape.left - REF_LEFT) < TOLERANCE:
                continue
            sp = shape.element
            sp.getparent().remove(sp)


def source_pictures(source_slide):
    pics = []
    for shape in source_slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if is_logo_pic(shape):
            continue
        if shape.width < 600000 or shape.height < 600000:
            continue
        pics.append(shape)
    return pics


def source_texts(source_slide, src_title):
    texts = []
    ts = get_title_shape(source_slide)
    for shape in source_slide.shapes:
        if not shape.has_text_frame or shape is ts:
            continue
        t = shape.text_frame.text.strip()
        if not t or t.split("\n")[0] == src_title:
            continue
        if is_logo_pic(shape):
            continue
        texts.append(t)
    texts.sort(key=lambda t: -len(t))
    return texts


def fill_slide_from_source(slide, title_shape, source_slide, new_title):
    set_title_style(title_shape, new_title)
    add_blue_divider(slide)
    clear_content_shapes(slide, title_shape)

    pics = source_pictures(source_slide)
    if len(pics) >= 2:
        for pic in pics:
            slide.shapes.add_picture(
                io.BytesIO(pic.image.blob), pic.left, pic.top, pic.width, pic.height
            )
    elif len(pics) == 1:
        pic = pics[0]
        left, top = 838200, 1000000
        max_w, max_h = 9600000, 5200000
        ratio = min(max_w / pic.width, max_h / pic.height, 1.0)
        w, h = int(pic.width * ratio), int(pic.height * ratio)
        slide.shapes.add_picture(io.BytesIO(pic.image.blob), left, top, w, h)
    else:
        texts = source_texts(source_slide, clean_title(get_title(source_slide)))
        body = "\n".join(texts[:6]) if texts else "（内容来自源课件，请结合图示讲解。）"
        tx = slide.shapes.add_textbox(838200, 1000000, 9600000, 5200000)
        set_body_text(tx, body[:2000], size_pt=18)


def resolve_source_path(work_dir, key, source_map):
    mapped = source_map.get(key, key)
    if os.path.isabs(mapped) and os.path.exists(mapped):
        return mapped
    candidate = os.path.join(work_dir, mapped)
    if os.path.exists(candidate):
        return candidate
    for path in glob.glob(os.path.join(work_dir, "*.pptx")):
        if mapped in os.path.basename(path) or key in os.path.basename(path):
            return path
    return None


def resolve_insert_index(prs, anchor):
    for i, slide in enumerate(prs.slides):
        t = clean_title(get_title(slide) or "")
        if norm(anchor) == norm(t):
            return i + 1
    for i, slide in enumerate(prs.slides):
        t = clean_title(get_title(slide) or "")
        if anchor in t:
            return i + 1
    return None


def make_prep_notes_200(title, slide, reason_prefix=""):
    title = clean_title(title)
    bits = []
    ts = get_title_shape(slide)
    for shape in slide.shapes:
        if not shape.has_text_frame or shape is ts:
            continue
        t = shape.text_frame.text.strip()
        if t and t.split("\n")[0] != title:
            bits.append(t[:100])
    preview = bits[0][:80] if bits else title

    if "案例" in title:
        p = f"本页案例聚焦「{title}」，用真实场景揭示项目/服务管理中的典型失误与可复制的改进路径"
        r = "案例教学能激活学员的经验记忆，比纯理论更易形成风险预警"
        e = f"案例要点：{preview}…；可追问学员「若您是项目经理，第一步会采取什么动作」"
        rp = "请带回岗位后对照自查：是否存在同类隐患，并制定一项可执行的预防或纠偏措施"
    elif "研讨" in title or "练习" in title:
        p = "本页为研讨/练习，目标是让学员产出与自身组织相关的结论"
        r = "被动听课难以转化；结构化研讨促成从「知道」到「用到」"
        e = "建议分组 8–10 分钟讨论，每组 1 分钟分享，讲师点评共性误区"
        rp = "请记录本组结论，作为 30-60-90 天改进计划的输入"
    elif title in ("Q&A",):
        p = "开放问答是知识内化的最后一环，也是课程价值从课堂走向岗位的关键接口"
        r = "学员提问往往指向真实痛点；现场回应可显著降低实践门槛"
        e = "常见问题如 PMO 落地、最小可行流程选型、变更审批设计等，结合白天案例给出建议"
        rp = "请鼓励具体提问——问题越具体，答案越可执行"
    elif "XXXX" in title:
        p = "本页为插入页空白模板，不在授课过程中展示"
        r = "模板页用于保持课件版式一致性"
        e = "无需讲解"
        rp = "跳过本页"
    else:
        p = f"本页核心观点：{title}，是课程大纲明确要求掌握的关键知识点"
        r = "该内容与信息化项目立项、交付或 IT 服务管理岗位的日常决策直接相关"
        e = f"页面要点：{preview}…；建议结合图示逐条讲解，并在每点后追问关联性"
        rp = f"请牢固记住：{title}——这是本段学习的核心记忆点"

    body = f"【P-观点】{p}。\n【R-理由】{r}。\n【E-例证】{e}。\n【P-重申】{rp}。"
    if reason_prefix:
        body = reason_prefix + body
    return body


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def insert_one(prs, template_idx, source_path, source_page, insert_idx, title, reason):
    source_prs = Presentation(source_path)
    source_slide = source_prs.slides[source_page - 1]
    work_title = title if title.endswith("（插）") else f"{title}（插）"

    duplicate_slide(prs, template_idx)
    new_idx = len(prs.slides) - 1
    move_slide(prs, new_idx, insert_idx)

    slide = prs.slides[insert_idx]
    title_shape = get_title_shape(slide)
    if not title_shape:
        title_shape = slide.shapes.add_textbox(REF_LEFT, REF_TOP, 7696200, 419100)

    fill_slide_from_source(slide, title_shape, source_slide, work_title)
    note = make_prep_notes_200(work_title, slide, f"插入理由：{reason}\n")
    set_notes(slide, note)
    return {"title": work_title, "page": insert_idx + 1}


def run_insert(prs, plan_path, work_dir):
    with open(plan_path, encoding="utf-8") as f:
        plan = json.load(f)
    source_map = plan.get("source_map", {})
    inserts = [x for x in plan.get("inserts", []) if x.get("mode", "script") == "script"]
    template_idx = find_template_index(prs, plan.get("template_marker", "XXXX（插）"))

    resolved = []
    for item in inserts:
        path = resolve_source_path(work_dir, item["source_key"], source_map)
        if not path:
            continue
        idx = resolve_insert_index(prs, item["anchor_title"])
        if idx is None:
            continue
        final = item.get("final_title") or clean_title(get_title(Presentation(path).slides[item["source_page"] - 1]))
        nt = final if final.endswith("（插）") else f"{final}（插）"
        if any(norm(get_title(s)) == norm(nt) for s in prs.slides):
            continue
        resolved.append((idx, item, path, final))

    resolved.sort(key=lambda x: -x[0])
    inserted = []
    for idx, item, path, final in resolved:
        tidx = find_template_index(prs)
        info = insert_one(
            prs, tidx, path, item["source_page"], idx, final, item.get("reason", "")
        )
        inserted.append(info)
    return inserted


def enrich_notes(prs, min_chars=200):
    updated = 0
    for i, slide in enumerate(prs.slides, 1):
        title = clean_title(get_title(slide) or "")
        if not title:
            for s in slide.shapes:
                if s.has_text_frame and "Q&A" in s.text_frame.text:
                    title = "Q&A"
                    break
            if not title and i == 1:
                title = "封面"

        existing = slide.notes_slide.notes_text_frame.text.strip()
        prefix = ""
        if existing.startswith("插入理由："):
            first_line = existing.split("\n")[0]
            rest = existing[len(first_line) :].strip()
            if len(rest) >= min_chars and rest.startswith("【P-观点】"):
                continue
            prefix = first_line + "\n"

        if len(existing) >= min_chars and existing.startswith("【P-观点】") and not prefix:
            continue

        note = make_prep_notes_200(title or f"第{i}页", slide, prefix)
        while len(note) < min_chars:
            note += "\n【延伸】请补充 1–2 个本单位实例，帮助学员完成从知识到行为的转化。"
        set_notes(slide, note)
        updated += 1
    return updated


def fix_dividers(prs):
    fixed = 0
    for slide in prs.slides:
        t = get_title(slide)
        if not t or "XXXX" in t:
            continue
        before = sum(
            1
            for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and s.width > 5000000
            and s.height < 80000
        )
        add_blue_divider(slide)
        after = sum(
            1
            for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and s.width > 5000000
            and s.height < 80000
        )
        if after > before:
            fixed += 1
    return fixed


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--ppt", required=True)
    parser.add_argument("--plan", default="insert_plan.json")
    parser.add_argument("--phase", choices=["insert", "notes", "fix-dividers", "all"], default="all")
    parser.add_argument("--min-chars", type=int, default=200)
    parser.add_argument("--work-dir", default=".")
    args = parser.parse_args()

    prs = Presentation(args.ppt)
    work_dir = os.path.dirname(os.path.abspath(args.ppt)) if args.work_dir == "." else args.work_dir

    if args.phase in ("insert", "all") and os.path.exists(args.plan):
        ins = run_insert(prs, args.plan, work_dir)
        print(f"Inserted {len(ins)} slides")
    if args.phase in ("fix-dividers", "all"):
        print(f"Fixed dividers on {fix_dividers(prs)} slides")
    if args.phase in ("notes", "all"):
        print(f"Updated notes on {enrich_notes(prs, args.min_chars)} slides")

    prs.save(args.ppt)
    print(f"Saved {args.ppt}")


if __name__ == "__main__":
    main()
