# -*- coding: utf-8 -*-
"""Index source PPT slides: title, pics, shapes, composite flag."""

import argparse
import glob
import hashlib
import json
import os
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

REF_LEFT, REF_TOP = 1409700, 152400
TOLERANCE = 500000


def get_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if abs(shape.left - REF_LEFT) < TOLERANCE and abs(shape.top - REF_TOP) < TOLERANCE:
                t = shape.text_frame.text.strip()
                if t:
                    return t.split("\n")[0]
    for shape in sorted(
        (s for s in slide.shapes if s.has_text_frame and s.top < 900000),
        key=lambda s: s.top,
    ):
        t = shape.text_frame.text.strip()
        if t and len(t.split("\n")[0]) < 100:
            return t.split("\n")[0]
    return ""


def count_pics(slide):
    n = 0
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if shape.left > 9000000 and shape.top < 500000:
            continue
        if shape.width < 600000 or shape.height < 600000:
            continue
        n += 1
    return n


def has_group_or_line(slide):
    for shape in slide.shapes:
        st = shape.shape_type
        if st in (MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
            return True
    return False


def index_file(path):
    prs = Presentation(path)
    rows = []
    for i, slide in enumerate(prs.slides, 1):
        sc = len(slide.shapes)
        pc = count_pics(slide)
        composite = sc >= 15 or has_group_or_line(slide)
        rows.append(
            {
                "file": os.path.basename(path),
                "page": i,
                "title": get_title(slide),
                "shape_count": sc,
                "pic_count": pc,
                "composite": composite,
            }
        )
    return rows


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("work_dir")
    parser.add_argument("--glob", default="*.pptx")
    args = parser.parse_args()

    paths = sorted(glob.glob(os.path.join(args.work_dir, args.glob)))
    all_rows = []
    for path in paths:
        if path.endswith("_work.pptx") or "temp" in os.path.basename(path).lower():
            continue
        try:
            all_rows.extend(index_file(path))
        except Exception as e:
            print(f"Skip {path}: {e}")

    out = os.path.join(args.work_dir, "source_slide_index.json")
    with open(out, "w", encoding="utf-8") as f:
        json.dump(all_rows, f, ensure_ascii=False, indent=2)
    composite = [r for r in all_rows if r["composite"] and r["title"]]
    print(f"Indexed {len(all_rows)} slides from {len(paths)} files")
    print(f"Composite candidates: {len(composite)}")
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
