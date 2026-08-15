# -*- coding: utf-8 -*-
"""
Generate 200+ char broadcast scripts (口播稿) for every slide in a PPT.

Usage:
  python broadcast_notes.py --input course.pptx --output course_v7.pptx
  python broadcast_notes.py --input course.pptx --output course_v7.pptx --min-chars 200 --course-name "培训名称"
"""

import argparse
import json
import os
import re
import shutil

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

REF_LEFT, REF_TOP = 1409700, 152400
TOLERANCE = 500000


def get_title_shape(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if abs(shape.left - REF_LEFT) < TOLERANCE and abs(shape.top - REF_TOP) < TOLERANCE:
            if shape.text_frame.text.strip():
                return shape
    for shape in sorted(
        (s for s in slide.shapes if s.has_text_frame and s.top < 900000),
        key=lambda s: s.top,
    ):
        t = shape.text_frame.text.strip()
        if t and len(t.split("\n")[0]) < 100:
            return shape
    return None


def get_title(slide):
    shape = get_title_shape(slide)
    return shape.text_frame.text.strip().split("\n")[0] if shape else ""


def is_logo(shape):
    return shape.left > 9000000 and shape.top < 500000


def extract_bodies(slide, title):
    ts = get_title_shape(slide)
    bodies = []
    for shape in slide.shapes:
        if not shape.has_text_frame or shape is ts:
            continue
        if is_logo(shape):
            continue
        t = shape.text_frame.text.strip()
        if not t or t.split("\n")[0] == title or len(t) < 4:
            continue
        bodies.append(t)
    bodies.sort(key=len, reverse=True)
    return bodies


def count_pics(slide):
    return sum(
        1
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
        and s.width > 800000
        and not is_logo(s)
    )


def snippet(text, n=100):
    t = re.sub(r"\s+", " ", text.replace("\x0b", " ")).strip()
    t = re.sub(r"^\d+[\.\、\s]+", "", t)
    return t[:n] if t else ""


def bullets(text, limit=3):
    lines = []
    for line in text.split("\n"):
        line = re.sub(r"^\d+[\.\、\s]+", "", line.strip())
        if len(line) >= 4:
            lines.append(line)
    return lines[:limit]


def pad_to_min(text, min_chars, extra=""):
    t = text.strip()
    fillers = [
        "讲完后可以停顿两秒，观察学员反应，再邀请一位同事分享他的理解。",
        "如果现场有不同做法，不必急于评判，先记录差异，课后可以整理成改进清单。",
        "这部分内容建议结合您本单位的真实场景来举例，会更有共鸣。",
        "请大家把这一页的要点记在工作手册上，回去对照岗位实践做一次自检。",
    ]
    i = 0
    while len(t) < min_chars and i < len(fillers):
        t += fillers[i]
        i += 1
    if extra and len(t) < min_chars:
        t += extra
    return t


def detect_role(page, total, title, bodies, pics, all_text):
    blob = title + all_text
    body_len = sum(len(b) for b in bodies)
    if page == 1:
        return "cover"
    if page == total:
        return "closing"
    if "Q&A" in blob or title.upper() == "Q&A":
        return "qa"
    if re.search(r"第[一二三四五六七八九十\d]+章", title) or "Chapter" in title:
        return "chapter"
    if "案例" in title:
        return "case"
    if "研讨" in title or "练习" in title:
        return "workshop"
    if body_len >= 80:
        return "concept"
    if pics >= 1:
        return "diagram"
    if "目录" in title:
        return "toc"
    return "concept"


def make_broadcast(page, total, title, bodies, pics, prev_title, course_name, min_chars):
    all_text = " ".join(bodies)
    role = detect_role(page, total, title, bodies, pics, all_text)
    topic = title or f"第{page}页"
    lead = snippet(bodies[0]) if bodies else ""
    bl = bullets(bodies[0]) if bodies else []

    if role == "cover":
        return pad_to_min(
            f"各位同事，大家好！欢迎参加今天的《{course_name}》。"
            "我是本次课程的主讲老师。接下来请大家准备好学员手册，跟着节奏走，"
            "中间会有案例和互动，建议随手记录。有疑问我们留到最后的问答环节。",
            min_chars,
        )
    if role == "closing":
        return pad_to_min(
            "好的，我们今天的课程到这里就接近尾声了。感谢大家一整天的专注投入。"
            "希望大家带着今天记录的改进点回到岗位，按30-60-90天的节奏逐步落地。"
            "也祝愿大家持续学习、不断精进，我们后会有期！",
            min_chars,
        )
    if role == "qa":
        return pad_to_min(
            "现在我们进入开放式问答环节。请大家把仍然困惑、或回到岗位立刻要解决的问题提出来。"
            "问题越具体，我们讨论越有价值。我会尽量结合今天的案例给出可执行的建议。",
            min_chars,
        )
    if role == "chapter":
        return pad_to_min(
            f"好，接下来我们进入新的模块：{topic}。"
            "切换话题前，先快速回顾上一模块的收获。进入本章后请调整注意力，跟着我的节奏走。",
            min_chars,
        )
    if role == "workshop":
        return pad_to_min(
            f"这一页是互动研讨，主题是「{topic}」。"
            f"建议分组讨论五到八分钟，围绕：{lead or '对照现状找差距、找一条可落地的改进动作'}。"
            "每组选一位代表准备一分钟发言，分享结束后我会做简短点评。",
            min_chars,
        )
    if role == "case":
        pts = "、".join(bl[:2]) if bl else "背景、问题、后果与启示"
        return pad_to_min(
            f"下面通过案例加深理解：「{topic}」。关键信息包括：{pts}。"
            "请对照自己：类似情况是否发生过？若您是负责人，第一步会做什么？",
            min_chars,
        )
    if role == "diagram":
        return pad_to_min(
            f"各位请看这一页的图，主题是「{topic}」。"
            f"请跟着讲解顺序：先看{snippet(lead, 60) or '整体框架'}，再看模块如何衔接。"
            "先抓住三个关键节点：输入、负责人、输出。对照本单位找最薄弱的一环。",
            min_chars,
        )
    if role == "toc":
        return pad_to_min(
            "这一页是课程目录。建议大家标出最关注的三个专题，课间可以重点提问。",
            min_chars,
        )
    if bl:
        joined = "；".join(bl[:3])
        script = (
            f"这一页我们讲「{topic}」。要点：{joined}。"
            "请边听边想：和您日常工作哪项直接相关？做不到位的后果是什么？"
            "回到岗位后可以立刻改的一个小动作是什么？"
        )
    elif lead:
        script = (
            f"这一页的主题是「{topic}」。核心信息：{lead}。"
            "这是课程的重要基础，后面很多内容都会用到。有不熟悉的名词先记下来。"
        )
    else:
        script = f"接下来看「{topic}」。这一页是关键知识点，请大家在手册上对应位置做标记。"

    if prev_title and prev_title == title and pics >= 1:
        script = f"我们继续看「{topic}」——这一页用图把刚才的概念具体化。" + script.split("。", 1)[-1]

    return pad_to_min(script, min_chars)


def run(input_path, output_path, min_chars, course_name, report_path=None):
    shutil.copy2(input_path, output_path)
    prs = Presentation(output_path)
    total = len(prs.slides)
    prev_title = ""
    pages = []

    for i, slide in enumerate(prs.slides, 1):
        title = get_title(slide)
        bodies = extract_bodies(slide, title)
        pics = count_pics(slide)
        note = make_broadcast(i, total, title, bodies, pics, prev_title, course_name, min_chars)
        slide.notes_slide.notes_text_frame.text = note
        pages.append({"page": i, "title": title, "notes_len": len(note)})
        if title:
            prev_title = title

    prs.save(output_path)

    lens = [p["notes_len"] for p in pages]
    summary = {
        "input": input_path,
        "output": output_path,
        "total": total,
        "notes_avg": round(sum(lens) / max(total, 1)),
        "notes_min": min(lens) if lens else 0,
        "notes_200plus": sum(1 for x in lens if x >= min_chars),
    }
    if report_path:
        with open(report_path, "w", encoding="utf-8") as f:
            json.dump({"summary": summary, "pages": pages}, f, ensure_ascii=False, indent=2)
    return summary


def main():
    parser = argparse.ArgumentParser(description="Generate broadcast notes for PPT slides")
    parser.add_argument("--input", "-i", required=True)
    parser.add_argument("--output", "-o", required=True)
    parser.add_argument("--min-chars", type=int, default=200)
    parser.add_argument("--course-name", default="培训课程")
    parser.add_argument("--report", default=None)
    args = parser.parse_args()

    s = run(args.input, args.output, args.min_chars, args.course_name, args.report)
    print(f"Saved: {s['output']}")
    print(f"Pages: {s['total']} | Avg: {s['notes_avg']} | Min: {s['notes_min']}")
    print(f">={args.min_chars} chars: {s['notes_200plus']}/{s['total']}")


if __name__ == "__main__":
    main()
