# -*- coding: utf-8 -*-
"""Verify PPT quality against V6 thresholds."""

import argparse
import json
import os
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

REF_LEFT, REF_TOP = 1409700, 152400
TOLERANCE = 500000


def get_title(slide):
    for s in slide.shapes:
        if s.has_text_frame and abs(s.left - REF_LEFT) < TOLERANCE and abs(s.top - REF_TOP) < TOLERANCE:
            t = s.text_frame.text.strip()
            if t:
                return t.split("\n")[0]
    return ""


def analyze(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        t = get_title(slide)
        notes = slide.notes_slide.notes_text_frame.text.strip()
        pics = sum(
            1
            for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE
            and s.width > 800000
            and not (s.left > 9000000 and s.top < 500000)
        )
        has_divider = any(
            s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and s.width > 5000000
            and s.height < 80000
            and 500000 < s.top < 800000
            for s in slide.shapes
        )
        slides.append(
            {
                "page": i,
                "title": t,
                "notes_len": len(notes),
                "pics": pics,
                "divider": has_divider,
                "prep": notes.startswith("【P-观点】") or notes.startswith("插入理由："),
            }
        )

    total = len(slides)
    pic_slides = sum(1 for s in slides if s["pics"] > 0)
    notes_200 = sum(1 for s in slides if s["notes_len"] >= 200)
    prep = sum(1 for s in slides if s["prep"])
    insert_marked = sum(1 for s in slides if "（插）" in s["title"] and "XXXX" not in s["title"])
    def exempt_short(s):
        t = s["title"]
        if not t or "XXXX" in t:
            return True
        if any(k in t for k in ("Chapter", "目录", "封面", "致谢")):
            return True
        if re.match(r"第[一二三四五六七八九十\d]+章", t):
            return True
        return False

    short = [s for s in slides if s["notes_len"] < 200 and not exempt_short(s)]

    titles = {}
    for s in slides:
        if s["title"]:
            titles.setdefault(s["title"], []).append(s["page"])
    dup_groups = {k: v for k, v in titles.items() if len(v) > 1}

    return {
        "path": path,
        "total_pages": total,
        "pic_slides": pic_slides,
        "notes_avg": round(sum(s["notes_len"] for s in slides) / max(total, 1)),
        "notes_200plus": notes_200,
        "prep_notes": prep,
        "insert_pages": insert_marked,
        "short_notes_pages": short,
        "duplicate_title_groups": len(dup_groups),
        "slides": slides,
    }


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("ppt")
    parser.add_argument("--gold", default=None)
    parser.add_argument("--min-notes", type=int, default=200)
    parser.add_argument("--report", default=None)
    args = parser.parse_args()

    result = analyze(args.ppt)
    total = result["total_pages"]
    pct_200 = result["notes_200plus"] / max(total, 1) * 100
    pct_prep = result["prep_notes"] / max(total, 1) * 100

    print(f"File: {args.ppt}")
    print(f"Pages: {total} | Pic slides: {result['pic_slides']} | Avg notes: {result['notes_avg']}")
    print(f"Notes>={args.min_notes}: {result['notes_200plus']} ({pct_200:.0f}%) | PREP: {result['prep_notes']} ({pct_prep:.0f}%)")
    print(f"Insert marked: {result['insert_pages']} | Dup title groups: {result['duplicate_title_groups']}")
    print(f"Short notes: {len(result['short_notes_pages'])}")

    pass_v6 = (
        230 <= total <= 260
        and result["pic_slides"] >= 90
        and pct_200 >= 85
        and pct_prep >= 80
    )
    print(f"V6 threshold: {'PASS' if pass_v6 else 'NEEDS WORK'}")

    if args.gold and os.path.exists(args.gold):
        gold = analyze(args.gold)
        print(f"\nGold {args.gold}: pages={gold['total_pages']} pics={gold['pic_slides']} notes200={gold['notes_200plus']}")

    if args.report:
        with open(args.report, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print(f"Report -> {args.report}")


if __name__ == "__main__":
    main()
