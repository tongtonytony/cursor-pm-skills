# -*- coding: utf-8 -*-
"""
PPT 内容与图片提取脚本（供 pptx-to-short-video 技能使用）

用法：python extract_pptx.py <PPT完整路径> [输出目录]
  输出目录默认：脚本所在技能的上级目录

输出：
  - pptx_extracted.txt：每页标题、要点、布局说明
  - img/：提取的图片（slide_N_img_J.扩展名）
"""

import sys
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
DEFAULT_OUT = SCRIPT_DIR.parent


def extract_pptx(pptx_path: Path, out_dir: Path) -> tuple[str, int]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(pptx_path))
    lines = []
    img_count = 0
    img_dir = out_dir / "img"
    img_dir.mkdir(parents=True, exist_ok=True)

    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n{'='*60}\n第 {idx} 页\n{'='*60}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
                lines.append("")
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        lines.append(t)
                lines.append("")
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append(" | ".join(cells))
                lines.append("")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    ext = (img.ext or "png").lstrip(".")
                    fname = f"slide_{idx:02d}_img_{img_count}.{ext}"
                    (img_dir / fname).write_bytes(img.blob)
                    lines.append(f"[图片: img/{fname}]")
                    img_count += 1
                except Exception:
                    pass
        lines.append("")

    return "\n".join(lines).strip(), img_count


def main():
    if len(sys.argv) < 2:
        print("用法: python extract_pptx.py <PPT完整路径> [输出目录]")
        sys.exit(1)
    ppt_path = Path(sys.argv[1])
    out_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else DEFAULT_OUT
    if not ppt_path.is_file():
        print(f"未找到文件: {ppt_path}")
        sys.exit(1)
    try:
        text, n_img = extract_pptx(ppt_path, out_dir)
    except ImportError:
        print("请先安装: pip install python-pptx")
        sys.exit(1)
    out_txt = out_dir / "pptx_extracted.txt"
    out_txt.write_text(text, encoding="utf-8")
    print(f"已生成: {out_txt}")
    print(f"图片: {out_dir / 'img'} (共 {n_img} 张)")


if __name__ == "__main__":
    main()
